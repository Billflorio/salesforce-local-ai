import sys
import os
import platform
import subprocess
import shutil
import tempfile
from pptx import Presentation

if hasattr(sys.stdout, 'reconfigure'):
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except Exception:
        pass

def find_piper():
    """Locate piper binary and an .onnx model file across Windows, Linux, and macOS."""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    system = platform.system()
    
    if system == "Windows":
        candidates = [
            os.path.join(script_dir, "piper", "piper", "piper.exe"),
            os.path.join(script_dir, "piper", "piper.exe"),
            os.path.join(script_dir, "piper.exe"),
            "piper.exe",
            "piper"
        ]
    else:
        candidates = [
            os.path.join(script_dir, "piper", "piper"),
            os.path.join(script_dir, "piper", "piper", "piper"),
            os.path.join(script_dir, "piper"),
            "piper"
        ]
    
    piper_bin = None
    for b in candidates:
        if os.path.isabs(b) or os.path.sep in b:
            if os.path.isfile(b) and (system == "Windows" or os.access(b, os.X_OK)):
                piper_bin = b
                break
        else:
            found = shutil.which(b)
            if found:
                piper_bin = found
                break
                
    if not piper_bin:
        return None, None
        
    search_dirs = [
        os.path.join(script_dir, "piper"),
        os.path.join(script_dir, "piper", "piper"),
        script_dir,
        os.path.dirname(piper_bin)
    ]
    model_path = None
    for s_dir in search_dirs:
        if os.path.isdir(s_dir):
            for fname in os.listdir(s_dir):
                if fname.endswith(".onnx"):
                    model_path = os.path.join(s_dir, fname)
                    break
        if model_path:
            break
            
    if piper_bin and model_path:
        return piper_bin, model_path
    return None, None

def play_wav(wav_path):
    system = platform.system()
    if system == "Windows":
        clean_path = wav_path.replace("'", "''")
        ps_cmd = f"(New-Object System.Media.SoundPlayer '{clean_path}').PlaySync()"
        subprocess.run(["powershell", "-Command", ps_cmd])
    elif system == "Darwin":
        subprocess.run(["afplay", wav_path])
    elif system == "Linux":
        if subprocess.run(["which", "paplay"], capture_output=True).returncode == 0:
            subprocess.run(["paplay", wav_path])
        elif subprocess.run(["which", "aplay"], capture_output=True).returncode == 0:
            subprocess.run(["aplay", wav_path])
        else:
            print(f"Generated {wav_path} but couldn't find paplay or aplay.")

def speak_text(text):
    system = platform.system()
    
    # 1. Try Piper neural TTS first on all platforms
    piper_bin, model_path = find_piper()
    if piper_bin and model_path:
        print(f"[PIPER TTS] Using Piper Neural Voice: {os.path.basename(model_path)}")
        wav_path = os.path.join(tempfile.gettempdir(), "slide_speech.wav")
        piper_proc = subprocess.Popen(
            [piper_bin, "--model", model_path, "--output_file", wav_path],
            stdin=subprocess.PIPE
        )
        piper_proc.communicate(input=text.encode("utf-8"))
        if os.path.exists(wav_path):
            play_wav(wav_path)
            return

    # 2. Platform-native fallback if Piper is not installed
    print(f"[SYSTEM TTS] Piper neural voice not found. Falling back to native {system} TTS...")
    if system == "Windows":
        clean_text = text.replace("'", "''").replace('"', '`"')
        ps_cmd = f"Add-Type -AssemblyName System.Speech; $synth = New-Object System.Speech.Synthesis.SpeechSynthesizer; $synth.Rate = 1; $synth.Speak('{clean_text}')"
        subprocess.run(["powershell", "-Command", ps_cmd])
    elif system == "Linux":
        if subprocess.run(["which", "spd-say"], capture_output=True).returncode == 0:
            subprocess.run(["spd-say", text])
        elif subprocess.run(["which", "espeak-ng"], capture_output=True).returncode == 0:
            subprocess.run(["espeak-ng", text])
        elif subprocess.run(["which", "espeak"], capture_output=True).returncode == 0:
            subprocess.run(["espeak", text])
        else:
            print("No Linux TTS engine found! Install piper (see PIPER_SETUP.md), espeak-ng, or speech-dispatcher (`sudo apt install espeak-ng`).")
    elif system == "Darwin":  # macOS
        subprocess.run(["say", text])
    else:
        print(f"Unsupported OS: {system}")

def main():
    if len(sys.argv) < 2:
        print("Usage: python speak_slide.py <slide_number_1_indexed> [optional_deck_path]")
        return
    
    slide_num = int(sys.argv[1])
    deck_path = sys.argv[2] if len(sys.argv) > 2 else 'Final_AIOnDimesDeck.pptx'
    
    if not os.path.exists(deck_path):
        # Fallback to absolute or other known names if relative fails
        script_dir = os.path.dirname(os.path.abspath(__file__))
        alt_path = os.path.join(script_dir, deck_path)
        if os.path.exists(alt_path):
            deck_path = alt_path
        else:
            print(f"Error: Presentation file not found at '{deck_path}'")
            return

    prs = Presentation(deck_path)
    
    if slide_num < 1 or slide_num > len(prs.slides):
        print(f"Error: Slide must be between 1 and {len(prs.slides)}")
        return
        
    slide = prs.slides[slide_num - 1]
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text
        print(f"\n--- SPEAKING SLIDE {slide_num} NOTES ---")
        print(notes)
        print("---------------------------------------\n")
        speak_text(notes)
    else:
        print(f"No notes found on slide {slide_num}")

if __name__ == '__main__':
    main()
